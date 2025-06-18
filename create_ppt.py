from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Helper to add image with title and explanation
def add_image_slide(prs, title, image_path, explanation=None):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(1)
    top = Inches(1.2)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(5.5))
    if explanation:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(8.5), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = explanation
        p.font.size = Pt(16)

# Create presentation
prs = Presentation()

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)

# 1. Title Slide
add_title_slide(prs, "Titanic Survival Prediction Project", "A Data Science Journey: EDA to Ensembling")

# 2. Agenda
add_bullet_slide(prs, "Agenda", [
    "Project Overview",
    "Exploratory Data Analysis (EDA)",
    "Data Cleaning & Feature Engineering",
    "Data Modeling",
    "Ensembling",
    "Results & Conclusions"
])

# 3. Project Overview
add_bullet_slide(prs, "Project Overview", [
    "Goal: Predict survival on the Titanic using machine learning.",
    "Dataset: Titanic passenger data (Kaggle)",
    "Process: EDA, feature engineering, modeling, and ensembling."
])

# 4. EDA - Add charts with explanations
add_bullet_slide(prs, "Exploratory Data Analysis (EDA)", [
    "Explored distributions and relationships between features and survival.",
    "Key findings: Age and Fare have wide distributions; survival rate ~38%."
])
add_image_slide(prs, "Ages Distribution", r"charts/ages_distribution.png", "Shows the age distribution of passengers. Most were between 20-40 years old, with a normal-like distribution.")
add_image_slide(prs, "Fare Distribution", r"charts/fare_distribution.png", "Fare values are right-skewed, with most passengers paying lower fares but a few paying much higher.")
add_image_slide(prs, "Survived People Percentages", r"charts/survived_people_precetanges.png", "Overall survival rate was about 38%. This chart visualizes the proportion of survivors vs. non-survivors.")
add_image_slide(prs, "Survived People Percentages in Each Relatives People Count", r"charts/survived_people_percentages_in_each relativesPeople_count.png", "Shows how survival rates varied depending on the number of relatives (siblings/spouses/parents/children) aboard.")

# 5. Data Cleaning & Feature Engineering
add_bullet_slide(prs, "Data Cleaning & Feature Engineering", [
    "Handled missing values and standardized formats.",
    "Engineered features: split Cabin into cell number and letter, extracted titles from Name, processed Ticket.",
    "Dropped irrelevant or redundant columns."
])

# 6. Data Modeling (with diagram)
add_bullet_slide(prs, "Data Modeling", [
    "Logistic Regression: This model looks for patterns in the data to predict who survived. It correctly predicted survival for 88 out of every 100 people.",
    "Random Forest: This is like asking a group of decision-makers to vote on each prediction. It got 88 out of 100 predictions right.",
    "Support Vector Machine: This model tries to draw the best possible line between those who survived and those who didn’t. It was correct 89% of the time.",
    "The F1 score (a single number that tells us how good the model is at both finding survivors and not making too many mistakes) was 87% for Logistic Regression, 87% for Random Forest, and 88% for SVM.",
    "Precision: Of all the people the model said would survive, how many actually did? (It’s about being careful not to give false hope.)",
    "Recall: Of all the people who really survived, how many did the model find? (It’s about not missing anyone.)",
    "The F1 score balances these two, so a high F1 means the model is good at both.",
    "All models were tested carefully to make sure the results are reliable."
])
add_image_slide(prs, "Random Forest Model Diagram", r"charts/RandomForestDiagram.png", "Visual representation of the Random Forest model structure used in the project.")

# 7. Ensembling
add_bullet_slide(prs, "Ensembling Approach", [
    "Combined predictions from Logistic Regression, Random Forest, and SVM.",
    "Applied model-specific probability thresholds for binarization.",
    "Final prediction by majority voting (at least 2/3 models predict survival).",
    "Ensembling improved robustness and overall accuracy."
])

# 8. Results & Conclusions
add_bullet_slide(prs, "Results & Conclusions", [
    "Ensembling improved prediction robustness and accuracy.",
    "Majority voting combined strengths of individual models.",
    "Project demonstrates a full data science workflow from EDA to deployment-ready predictions."
])

prs.save(r"Titanic_DataScience_Project_Presentation.pptx")
